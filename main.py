import os
import tempfile
import json
import re

from flask import Flask, request, jsonify, abort
from flask_cors import CORS
from werkzeug.utils import secure_filename

import PyPDF2
from groq import Groq

import pytesseract
from pdf2image import convert_from_path
from PIL import Image, ImageFilter
import cv2
import numpy as np

from docx import Document
from pptx import Presentation

from dotenv import load_dotenv

# Load environment variables
load_dotenv()

app = Flask(__name__)
CORS(app, origins=["http://localhost:5173"], supports_credentials=True)

try:
    groq_api_key = os.getenv("GROQ_API_KEY")
    if not groq_api_key:
        raise ValueError("GROQ_API_KEY not found in environment variables")
    groq_client = Groq(api_key=groq_api_key)
except Exception as e:
    print(f"Error initializing Groq client: {e}")
    groq_client = None

SUPPORTED_EXTENSIONS = {
    ".pdf": "application/pdf",
    ".png": "image/png",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation"
}

MAX_FILE_SIZE = 10 * 1024 * 1024  # 10MB

def allowed_file(filename):
    ext = os.path.splitext(filename)[1].lower()
    return ext in SUPPORTED_EXTENSIONS

def validate_file(file_storage):
    filename = file_storage.filename
    if not filename:
        return False, "No filename provided"
    if not allowed_file(filename):
        return False, f"Unsupported file type: {os.path.splitext(filename)[1].lower()}"
    file_storage.seek(0, os.SEEK_END)
    size = file_storage.tell()
    file_storage.seek(0)
    if size == 0:
        return False, "File is empty"
    if size > MAX_FILE_SIZE:
        return False, "File size too large (max 10MB)"
    return True, None

def preprocess_image_pil(image: Image.Image) -> Image.Image:
    image = image.convert("L")
    image = image.point(lambda x: 0 if x < 140 else 255, '1')
    image = image.filter(ImageFilter.SHARPEN)
    return image

def preprocess_image_cv2(np_img: np.ndarray) -> np.ndarray:
    gray = cv2.cvtColor(np_img, cv2.COLOR_BGR2GRAY)
    norm_img = np.zeros((gray.shape[0], gray.shape[1]))
    img = cv2.normalize(gray, norm_img, 0, 255, cv2.NORM_MINMAX)
    img = cv2.threshold(img, 100, 255, cv2.THRESH_BINARY)[1]
    img = cv2.GaussianBlur(img, (1, 1), 0)
    return img

def ocr_image_file(image_path: str) -> str:
    try:
        image = Image.open(image_path)
        image = preprocess_image_pil(image)
        text = pytesseract.image_to_string(image)
        if text.strip():
            return text
        np_img = cv2.imread(image_path)
        img = preprocess_image_cv2(np_img)
        text = pytesseract.image_to_string(img)
        return text
    except Exception as e:
        print(f"OCR image extraction failed: {e}")
        return ""

def ocr_extract_text_from_pdf(pdf_path: str) -> str:
    try:
        pages = convert_from_path(pdf_path, dpi=300)
        texts = []
        for page in pages:
            pre_img = preprocess_image_pil(page)
            page_text = pytesseract.image_to_string(pre_img)
            texts.append(page_text)
        return "\n".join(texts).strip()
    except Exception as e:
        print(f"OCR extraction failed: {e}")
        return ""

def extract_text_from_pdf_optimized(file_path: str) -> str:
    text = ""
    try:
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
        text = text.strip()
        if text:
            return text
    except Exception as e:
        print(f"Error reading PDF normally: {e}")
    ocr_text = ocr_extract_text_from_pdf(file_path)
    if ocr_text:
        return ocr_text
    else:
        return ""

def extract_text_from_docx(file_path: str) -> str:
    try:
        doc = Document(file_path)
        full_text = []
        for paragraph in doc.paragraphs:
            full_text.append(paragraph.text)
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    full_text.append(cell.text)
        return '\n'.join(full_text).strip()
    except Exception as e:
        print(f"Error extracting text from DOCX: {e}")
        return ""

def extract_text_from_pptx(file_path: str) -> str:
    try:
        prs = Presentation(file_path)
        full_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    full_text.append(shape.text)
        return '\n'.join(full_text).strip()
    except Exception as e:
        print(f"Error extracting text from PPTX: {e}")
        return ""

def extract_text_from_file(file_path: str, file_extension: str) -> str:
    if file_extension == ".pdf":
        return extract_text_from_pdf_optimized(file_path)
    elif file_extension == ".docx":
        return extract_text_from_docx(file_path)
    elif file_extension == ".pptx":
        return extract_text_from_pptx(file_path)
    elif file_extension in [".png", ".jpg", ".jpeg"]:
        return ocr_image_file(file_path)
    else:
        return ""

def extract_json_from_response(response_text: str) -> dict:
    match = re.search(r"\{.*\}", response_text, re.DOTALL)
    if match:
        json_str = match.group(0)
        return json.loads(json_str)
    return json.loads(response_text)

def clean_bullet_points(text):
    lines = text.splitlines()
    cleaned = []
    prev_line_bullet = False
    for line in lines:
        stripped = line.lstrip()
        if stripped.startswith(("-", "*", "•")):
            cleaned.append(stripped)
            prev_line_bullet = True
        elif prev_line_bullet and not stripped:
            continue
        else:
            if prev_line_bullet and cleaned:
                cleaned[-1] += " " + stripped
            else:
                cleaned.append(stripped)
            prev_line_bullet = False
    return "\n".join(cleaned)

def classify_and_summarize_document(text: str, filename: str) -> dict:
    if not groq_client:
        abort(500, description="Groq client not initialized. Please check your API key.")
    cleaned_text = re.sub(r'\s+', ' ', text).strip()
    limited_text = cleaned_text[:3000] if len(cleaned_text) > 3000 else cleaned_text
    prompt = f"""
    Analyze the following document text from file "{filename}" and provide:
    1. Document Type: Classify this document as one of these types:
       - Offer Letter
       - Experience Letter
       - Medical Certificate
       - Leave Letter
       - Resignation Letter
       - Salary Certificate
       - Presentation/Slides (for PowerPoint files)
       - Report/Document (for Word documents)
       - Other (specify what type)
    2. Summary: Provide a one-line summary of the document's main content
    3. Confidence: Rate your confidence in the classification (High/Medium/Low)
    Please respond ONLY in valid JSON format with keys: "document_type", "summary", "confidence"
    Document Text:
    {limited_text}
    """
    try:
        response = groq_client.chat.completions.create(
            model="llama3-70b-8192",
            messages=[
                {
                    "role": "system",
                    "content": "You are an expert document classifier. Analyze documents and provide accurate classifications and summaries in valid JSON format only."
                },
                {
                    "role": "user",
                    "content": prompt
                }
            ],
            temperature=0.1,
            max_tokens=500
        )
        response_text = response.choices[0].message.content.strip()
        try:
            result = extract_json_from_response(response_text)
            return {
                "document_type": result.get("document_type", "Unknown"),
                "summary": result.get("summary", "Unable to generate summary"),
                "confidence": result.get("confidence", "Low")
            }
        except Exception as json_error:
            print(f"JSON parsing error: {json_error}")
            print(f"Response text: {response_text}")
            doc_type = "Unknown"
            summary = "Unable to generate summary"
            confidence = "Low"
            type_match = re.search(r'"?document_type"?\s*:\s*"([^"]*)"', response_text, re.IGNORECASE)
            summary_match = re.search(r'"?summary"?\s*:\s*"([^"]*)"', response_text, re.IGNORECASE)
            confidence_match = re.search(r'"?confidence"?\s*:\s*"([^"]*)"', response_text, re.IGNORECASE)
            if type_match:
                doc_type = type_match.group(1)
            if summary_match:
                summary = summary_match.group(1)
            if confidence_match:
                confidence = confidence_match.group(1)
            return {
                "document_type": doc_type,
                "summary": summary,
                "confidence": confidence
            }
    except Exception as e:
        print(f"Groq API error: {e}")
        abort(500, description=f"Error with Groq API: {str(e)}")

def parse_experience_years(text):
    matches = re.findall(r'(\d+)\s*(?:\+)?\s*(?:years?|yrs?)', text, re.IGNORECASE)
    years = [int(m) for m in matches]
    return max(years) if years else None

def get_experience_range(jd_text):
    match = re.search(r'(\d+)\s*-\s*(\d+)\s*(?:years?|yrs?)', jd_text, re.IGNORECASE)
    if match:
        return int(match.group(1)), int(match.group(2))
    match = re.search(r'(?:minimum|required)\s*(\d+)\s*(?:years?|yrs?)', jd_text, re.IGNORECASE)
    if match:
        return int(match.group(1)), None
    return None, None

def llm_score_education(jd_text: str, resume_text: str, groq_client) -> dict:
    # Always keep JD at the top, and dynamically adjust resume text slice
    jd_slice = jd_text[:2000]
    max_prompt_length = 4000  # adjust if needed
    available_for_resume = max_prompt_length - len(jd_slice)
    resume_slice = resume_text[:available_for_resume]

    prompt = f"""
You are an expert HR AI. Given the following job description and candidate resume, extract and compare the education requirements and qualifications.

Instructions:
1. Extract the required education level from the job description (e.g., "Bachelor's in Computer Science", "Master's degree", "PhD", etc.).
2. Extract the highest education level from the resume.
3. If the candidate's education meets or exceeds the requirement, return 50 for "education_score". Otherwise, return 0.
4. Provide a brief reason for the score.

Return ONLY valid JSON with these keys:
- "education_score": int (0 or 50)
- "education_reason": string
- "jd_education": string
- "resume_education": string

--- JOB DESCRIPTION ---
{jd_slice}

--- RESUME ---
{resume_slice}
"""
    # Debug: print prompt length and key sections
    print(f"JD chars: {len(jd_slice)}, Resume chars: {len(resume_slice)}, Total prompt: {len(prompt)}")

    try:
        response = groq_client.chat.completions.create(
            model="llama3-70b-8192",
            messages=[
                {
                    "role": "system",
                    "content": "You are an expert HR AI that returns only valid JSON with detailed scoring and reasoning."
                },
                {
                    "role": "user",
                    "content": prompt
                }
            ],
            temperature=0.1,
            max_tokens=500
        )
        response_text = response.choices[0].message.content.strip()
        return extract_json_from_response(response_text)
    except Exception as e:
        print(f"Groq API error during education scoring: {e}")
        return {
            "education_score": 0,
            "education_reason": "Error occurred during education scoring.",
            "jd_education": "",
            "resume_education": ""
        }

def score_resume_against_jd(jd_text: str, resume_text: str, resume_filename: str) -> dict:
    if not groq_client:
        abort(500, description="Groq client not initialized. Please check your API key.")

    # --- Experience scoring (out of 50) ---
    jd_min, jd_max = get_experience_range(jd_text)
    candidate_exp = parse_experience_years(resume_text)
    exp_score = 0
    exp_reason = ""
    max_exp_score = 50

    if candidate_exp is not None and jd_min is not None:
        if jd_max is not None:
            if jd_min <= candidate_exp <= jd_max:
                exp_score = max_exp_score
                exp_reason = f"Candidate's experience ({candidate_exp} years) is within the required range ({jd_min}-{jd_max} years)."
            else:
                exp_score = 0
                exp_reason = f"Candidate's experience ({candidate_exp} years) is outside the required range ({jd_min}-{jd_max} years)."
        else:
            if candidate_exp >= jd_min:
                exp_score = max_exp_score
                exp_reason = f"Candidate's experience ({candidate_exp} years) meets or exceeds the minimum required ({jd_min} years)."
            else:
                exp_score = 0
                exp_reason = f"Candidate's experience ({candidate_exp} years) is below the required minimum ({jd_min} years)."
    else:
        exp_score = 0
        exp_reason = "Could not determine experience from documents."

    # --- Education scoring (LLM-based, full or zero out of 50) ---
    edu_result = llm_score_education(jd_text, resume_text, groq_client)
    edu_score = edu_result.get("education_score", 0)
    edu_reason = edu_result.get("education_reason", "")

    # --- Skills scoring (out of 50) ---
    cleaned_resume_text = clean_bullet_points(resume_text)
    prompt = f"""
You are an expert HR AI. Given the following job description and candidate resume, perform a detailed evaluation:

1. **Skills Match (0-50):**
   - Score based on overlap between required and candidate skills.
   - Give full credit for exact matches, partial credit for closely related skills (e.g., Java vs. JavaScript), and explain your reasoning.
   - List skills as: exact_matches, partial_matches (with explanation), missing_skills.

2. **Feedback:**
   - Summarize the candidate's strengths and weaknesses for this role.
   - Suggest areas for improvement.

Return your answer as valid JSON with these keys (all are required, do not omit any!):
- "skills_score": int (0-50)
- "skills_reason": string
- "exact_matches": list of strings
- "partial_matches": list of dicts with "skill" and "reason"
- "missing_skills": list of strings
- "strengths": string
- "weaknesses": string

--- JOB DESCRIPTION ---
{jd_text[:2500]}

--- RESUME ---
{cleaned_resume_text[:2500]}
"""
    try:
        response = groq_client.chat.completions.create(
            model="llama3-70b-8192",
            messages=[
                {
                    "role": "system",
                    "content": "You are an expert HR AI that returns only valid JSON with detailed scoring and reasoning."
                },
                {
                    "role": "user",
                    "content": prompt
                }
            ],
            temperature=0.1,
            max_tokens=900
        )
        response_text = response.choices[0].message.content.strip()
        result = extract_json_from_response(response_text)
        skills_score = result.get("skills_score", 0)
        skills_reason = result.get("skills_reason", "")

        total_score = exp_score + skills_score + edu_score  # Out of 150
        return {
            "score": total_score,
            "experience_score": exp_score,
            "skills_score": skills_score,
            "education_score": edu_score,
            "qualification_status": "Qualified" if exp_score >= 30 and skills_score >= 30 and edu_score == 50 else "Underqualified",
            "summary": f"Experience: {exp_reason} | Skills: {skills_reason} | Education: {edu_reason}",
            "experience_match": exp_score >= 30,
            "experience_comment": exp_reason,
            "education_comment": edu_reason,
            "skills_matched": ", ".join(result.get("exact_matches", [])),
            "skills_partial": ", ".join([f"{pm['skill']} ({pm['reason']})" for pm in result.get("partial_matches", [])]),
            "skills_missing": ", ".join(result.get("missing_skills", [])),
            "strengths": result.get("strengths", ""),
            "weaknesses": result.get("weaknesses", "")
        }
    except Exception as e:
        print(f"Groq LLM scoring error: {e}")
        abort(500, description=f"Error with Groq API or LLM response: {str(e)}")

@app.route("/score-resumes", methods=["POST"])
def score_resumes():
    jd_file = request.files.get("jd_file")
    jd_text = request.form.get("jd_text", "").strip()
    resumes = request.files.getlist("resumes")

    if not jd_file and not jd_text:
        return jsonify({"detail": "No job description provided (file or text required)"}), 400

    if not resumes:
        return jsonify({"detail": "No resumes uploaded"}), 400

    jd_text_extracted = ""
    temp_jd_path = None

    try:
        # --- Extract JD text ONCE and use for all resumes ---
        if jd_file:
            valid, err = validate_file(jd_file)
            if not valid:
                return jsonify({"detail": f"JD file error: {err}"}), 400
            filename = secure_filename(jd_file.filename)
            ext = os.path.splitext(filename)[1].lower()
            if ext not in SUPPORTED_EXTENSIONS:
                return jsonify({"detail": f"Unsupported JD file type: {ext}"}), 400
            with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as temp_file:
                jd_file.save(temp_file.name)
                temp_jd_path = temp_file.name
            jd_text_extracted = extract_text_from_file(temp_jd_path, ext)
        else:
            jd_text_extracted = jd_text

        if not jd_text_extracted or len(jd_text_extracted) < 10:
            return jsonify({"detail": "No usable text in job description"}), 400

        # DEBUG: Log JD text for all resumes
        print("==== JD TEXT USED FOR ALL RESUMES ====")
        print(jd_text_extracted[:1000])  # Print first 1000 chars for debugging

        results = []
        failed_files = []

        for resume in resumes:
            if not resume.filename:
                failed_files.append("Unknown filename - No filename provided")
                continue
            valid, err = validate_file(resume)
            if not valid:
                failed_files.append(f"{resume.filename} - {err}")
                continue
            filename = secure_filename(resume.filename)
            ext = os.path.splitext(filename)[1].lower()
            temp_resume_path = None
            try:
                with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as temp_file:
                    resume.save(temp_file.name)
                    temp_resume_path = temp_file.name
                resume_text = extract_text_from_file(temp_resume_path, ext)
                if not resume_text or len(resume_text.strip()) < 10:
                    failed_files.append(f"{resume.filename} - No readable text extracted")
                    continue

                print(f"Scoring resume: {resume.filename} | Resume text length: {len(resume_text)}")

                result = score_resume_against_jd(jd_text_extracted, resume_text, resume.filename)
                results.append({
                    "filename": resume.filename,
                    "score": result["score"],
                    "experience_score": result["experience_score"],
                    "skills_score": result["skills_score"],
                    "education_score": result["education_score"],
                    "qualification_status": result["qualification_status"],
                    "summary": result["summary"],
                    "experience_match": result["experience_match"],
                    "experience_comment": result["experience_comment"],
                    "education_comment": result["education_comment"],
                    "skills_matched": result["skills_matched"],
                    "skills_partial": result["skills_partial"],
                    "skills_missing": result["skills_missing"],
                    "strengths": result["strengths"],
                    "weaknesses": result["weaknesses"]
                })
            except Exception as e:
                failed_files.append(f"{resume.filename} - Error: {str(e)}")
            finally:
                if temp_resume_path and os.path.exists(temp_resume_path):
                    try:
                        os.unlink(temp_resume_path)
                    except Exception:
                        pass

        return jsonify({
            "job_description": jd_text_extracted[:2000],
            "results": results,
            "failed_files": failed_files
        })
    finally:
        if temp_jd_path and os.path.exists(temp_jd_path):
            try:
                os.unlink(temp_jd_path)
            except Exception:
                pass
        jd_text_extracted = ""
        temp_jd_path = None

@app.route("/analyze-file", methods=["POST"])
def analyze_single_file():
    if "file" not in request.files:
        return jsonify({"detail": "No file uploaded"}), 400
    file = request.files["file"]
    valid, err = validate_file(file)
    if not valid:
        return jsonify({"detail": err}), 400

    filename = secure_filename(file.filename)
    ext = os.path.splitext(filename)[1].lower()
    temp_file_path = None

    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as temp_file:
            file.save(temp_file.name)
            temp_file_path = temp_file.name
        extracted_text = extract_text_from_file(temp_file_path, ext)
        if not extracted_text.strip():
            return jsonify({"detail": "No text could be extracted from the file."}), 400
        if len(extracted_text.strip()) < 10:
            return jsonify({"detail": "Insufficient text content in the file for analysis."}), 400
        analysis_result = classify_and_summarize_document(extracted_text, filename)
        return jsonify({
            "filename": filename,
            "document_type": analysis_result["document_type"],
            "summary": analysis_result["summary"],
            "confidence": analysis_result["confidence"],
            "extracted_text_length": len(extracted_text)
        })
    except Exception as e:
        print(f"Unexpected error in analyze_single_file: {e}")
        return jsonify({"detail": f"An unexpected error occurred: {str(e)}"}), 500
    finally:
        if temp_file_path and os.path.exists(temp_file_path):
            try:
                os.unlink(temp_file_path)
            except Exception:
                pass

@app.route("/analyze-multiple-files", methods=["POST"])
def analyze_multiple_files():
    files = request.files.getlist("files")
    if not files:
        return jsonify({"detail": "No files uploaded"}), 400
    if len(files) > 10:
        return jsonify({"detail": "Too many files. Maximum 10 files per request."}), 400

    analyses = []
    failed_files = []
    successful_count = 0

    for file in files:
        if not file.filename:
            failed_files.append("Unknown filename - No filename provided")
            continue
        valid, err = validate_file(file)
        if not valid:
            failed_files.append(f"{file.filename} - {err}")
            continue
        filename = secure_filename(file.filename)
        ext = os.path.splitext(filename)[1].lower()
        temp_file_path = None
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as temp_file:
                file.save(temp_file.name)
                temp_file_path = temp_file.name
            extracted_text = extract_text_from_file(temp_file_path, ext)
            if not extracted_text.strip():
                failed_files.append(f"{file.filename} - No text could be extracted")
                continue
            if len(extracted_text.strip()) < 10:
                failed_files.append(f"{file.filename} - Insufficient text content")
                continue
            analysis_result = classify_and_summarize_document(extracted_text, filename)
            analyses.append({
                "filename": filename,
                "document_type": analysis_result["document_type"],
                "summary": analysis_result["summary"],
                "confidence": analysis_result["confidence"],
                "extracted_text_length": len(extracted_text)
            })
            successful_count += 1
        except Exception as e:
            failed_files.append(f"{file.filename} - Processing error: {str(e)}")
        finally:
            if temp_file_path and os.path.exists(temp_file_path):
                try:
                    os.unlink(temp_file_path)
                except Exception:
                    pass

    return jsonify({
        "total_files": len(files),
        "successful_analyses": successful_count,
        "failed_files": failed_files,
        "analyses": analyses
    })

@app.route("/health", methods=["GET"])
def health_check():
    groq_status = "connected" if groq_client else "not connected"
    api_key_status = "present" if os.getenv("GROQ_API_KEY") else "missing"
    return jsonify({
        "status": "healthy",
        "message": "Document Analyzer API is running",
        "groq_status": groq_status,
        "api_key_status": api_key_status,
        "supported_file_types": list(SUPPORTED_EXTENSIONS.keys())
    })

@app.route("/", methods=["GET"])
def root():
    return jsonify({
        "message": "Document & Image Analyzer API",
        "version": "2.4.0",
        "supported_file_types": list(SUPPORTED_EXTENSIONS.keys()),
        "endpoints": {
            "analyze_single": "/analyze-file",
            "analyze_multiple": "/analyze-multiple-files",
            "score_resumes": "/score-resumes",
            "health": "/health"
        }
    })

if __name__ == "__main__":
    app.run(host="0.0.0.0", port=3000, debug=True) 
